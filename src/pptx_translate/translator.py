from __future__ import annotations

import json
import logging
import os
import secrets
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Dict, Iterable, List, Optional, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from pptx_translate.backends import TranslationBackend
from pptx_translate.models import (
    OcrImageInput,
    OcrTextRegion,
    TranslatableUnit,
)
from pptx_translate.ocr import OcrBackend, PytesseractOcrBackend

NAMESPACES = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "cp": "http://schemas.openxmlformats.org/officeDocument/2006/custom-properties",
    "vt": "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes",
}


@dataclass
class DocumentPart:
    """
    Holds a parsed XML part and its translatable text nodes.
    """

    path: Path
    tree: ET.ElementTree
    nodes: List[Tuple[ET.Element, TranslatableUnit]]


@dataclass
class DeckProfile:
    """
    Derived context from the deck to guide translation.
    """

    title: str | None
    section_titles: List[str]
    keywords: List[str]
    summary: str | None

    def to_context_string(self) -> str:
        parts: List[str] = []
        if self.title:
            parts.append(f"Deck title: {self.title}")
        if self.section_titles:
            parts.append("Sections: " + "; ".join(self.section_titles[:10]))
        if self.keywords:
            parts.append("Frequent terms: " + ", ".join(self.keywords[:15]))
        if self.summary:
            parts.append(f"Summary: {self.summary}")
        return "\n".join(parts)


class PptxTranslator:
    """
    Translate PowerPoint files while preserving layout/formatting.
    Supports optional OCR for text inside images, run metadata, and QA reporting.
    """

    def __init__(
        self,
        backend: TranslationBackend,
        include_notes: bool = True,
        include_masters: bool = True,
        max_batch_chars: int = 4000,
        dry_run: bool = False,
        dedupe_text: bool = True,
        translate_images: bool = False,
        ocr_backend: Optional[OcrBackend] = None,
        ocr_config: Optional[dict] = None,
        max_concurrent_requests: int = 1,
        profile: str = "balanced",
        qa_report_path: Optional[Path] = None,
        qa_report_format: str = "json",
        qa_threshold_length_ratio: float = 1.6,
    ) -> None:
        self.backend = backend
        self.include_notes = include_notes
        self.include_masters = include_masters
        self.max_batch_chars = max_batch_chars
        self.dry_run = dry_run
        self.dedupe_text = dedupe_text
        self.translate_images = translate_images
        self.ocr_backend = ocr_backend
        self.ocr_config = ocr_config or {}
        self.max_concurrent_requests = max_concurrent_requests
        self.profile = profile
        self.qa_report_path = qa_report_path
        self.qa_report_format = qa_report_format
        self.qa_threshold_length_ratio = qa_threshold_length_ratio
        self.logger = logging.getLogger(__name__)
        self._id_counter = 0

    def translate_file(
        self,
        input_path: Path,
        output_path: Path,
        source_lang: str | None,
        target_lang: str,
        glossary: list[dict] | None = None,
        context: str | None = None,
        run_id: Optional[str] = None,
        generate_glossary_path: Optional[Path] = None,
        deck_profile_path: Optional[Path] = None,
    ) -> List[TranslatableUnit]:
        """
        Main entrypoint: translate `input_path` into `output_path`.
        Returns list of translated units.
        """
        input_path = input_path.resolve()
        output_path = output_path.resolve()
        self._validate_input(input_path)
        run_id = run_id or generate_run_id()

        deck_profile = self._build_deck_profile(input_path)
        if deck_profile_path:
            deck_profile_path.write_text(deck_profile.to_context_string(), encoding="utf-8")

        combined_context = self._combine_context(context, deck_profile)

        with TemporaryDirectory() as tmp_dir:
            temp_root = Path(tmp_dir)
            self._unpack_pptx(input_path, temp_root)
            parts = self._load_parts(temp_root)
            text_units = [unit for part in parts for _, unit in part.nodes]

            if generate_glossary_path:
                self._generate_glossary(generate_glossary_path, text_units, deck_profile)
                self.logger.info("Generated glossary at %s", generate_glossary_path)
                return []

            ocr_units: List[TranslatableUnit] = []
            ocr_regions: List[OcrTextRegion] = []
            if self.translate_images:
                ocr_units, ocr_regions = self._extract_image_units(input_path)
                self.logger.info("Extracted %d OCR text regions from images", len(ocr_units))

            all_units = text_units + ocr_units
            if not all_units:
                self.logger.info("No translatable content found")
                return []

            translated_units = self._translate_units(
                all_units,
                source_lang=source_lang,
                target_lang=target_lang,
                glossary=glossary,
                context=combined_context,
            )
            translated_map: Dict[str, TranslatableUnit] = {u.id: u for u in translated_units}

            if not self.dry_run:
                self._inject_translations(parts, translated_map)
                self._repack_pptx(temp_root, output_path)
                if self.translate_images and ocr_regions:
                    self._apply_image_overlays(output_path, ocr_regions, translated_map)
                self._embed_run_metadata(
                    output_path,
                    {
                        "run_id": run_id,
                        "source_lang": source_lang or "auto",
                        "target_lang": target_lang,
                        "backend": self.backend.__class__.__name__,
                        "profile": self.profile,
                        "timestamp_utc": datetime.now(timezone.utc).isoformat(),
                    },
                )
                self.logger.info("Wrote translated file to %s (run_id=%s)", output_path, run_id)
                if self.qa_report_path:
                    self._generate_qa_report(output_path, translated_units, run_id)
            else:
                self.logger.info("Dry run mode: no output file written")

            return translated_units

    def _validate_input(self, input_path: Path) -> None:
        if not input_path.exists():
            raise FileNotFoundError(f"Input file not found: {input_path}")
        if input_path.suffix.lower() != ".pptx":
            raise ValueError("Input file must be a .pptx")

    def _unpack_pptx(self, input_path: Path, temp_root: Path) -> None:
        with zipfile.ZipFile(input_path, "r") as zf:
            zf.extractall(temp_root)
        self.logger.debug("Extracted pptx to %s", temp_root)

    def _load_parts(self, temp_root: Path) -> List[DocumentPart]:
        paths = self._discover_xml_parts(temp_root)
        parts: List[DocumentPart] = []

        for path in paths:
            tree = ET.parse(path)
            root = tree.getroot()
            nodes: List[Tuple[ET.Element, TranslatableUnit]] = []
            for idx, elem in enumerate(root.findall(".//a:t", namespaces=NAMESPACES)):
                text = elem.text if elem.text is not None else ""
                if text == "":
                    continue
                unit_id = self._next_id()
                location = f"{path.relative_to(temp_root)}::a:t[{idx}]"
                unit = TranslatableUnit(
                    id=unit_id,
                    location=str(location),
                    source_text=text,
                    context=None,
                )
                nodes.append((elem, unit))
            parts.append(DocumentPart(path=path, tree=tree, nodes=nodes))

        return parts

    def _translate_units(
        self,
        units: List[TranslatableUnit],
        source_lang: str | None,
        target_lang: str,
        glossary: list[dict] | None,
        context: str | None,
    ) -> List[TranslatableUnit]:
        if not self.dedupe_text:
            return self.backend.translate(
                units,
                source_lang=source_lang,
                target_lang=target_lang,
                max_batch_chars=self.max_batch_chars,
                glossary=glossary,
                context=context,
                max_concurrent_requests=self.max_concurrent_requests,
            )

        text_to_units: Dict[str, List[TranslatableUnit]] = {}
        unique_units: List[TranslatableUnit] = []
        for unit in units:
            key = unit.source_text
            if key not in text_to_units:
                text_to_units[key] = [unit]
                unique_units.append(unit)
            else:
                text_to_units[key].append(unit)

        self.logger.info("Deduped %d texts down to %d unique entries", len(units), len(unique_units))

        translated_unique = self.backend.translate(
            unique_units,
            source_lang=source_lang,
            target_lang=target_lang,
            max_batch_chars=self.max_batch_chars,
            glossary=glossary,
            context=context,
            max_concurrent_requests=self.max_concurrent_requests,
        )
        by_text: Dict[str, str] = {u.source_text: (u.translated_text or u.source_text) for u in translated_unique}

        translated_all: List[TranslatableUnit] = []
        for unit in units:
            translated_text = by_text.get(unit.source_text, unit.source_text)
            translated_all.append(
                TranslatableUnit(
                    id=unit.id,
                    location=unit.location,
                    source_text=unit.source_text,
                    translated_text=translated_text,
                    context=unit.context,
                )
            )
        return translated_all

    def _inject_translations(
        self,
        parts: List[DocumentPart],
        translated_map: Dict[str, TranslatableUnit],
    ) -> None:
        for part in parts:
            for elem, unit in part.nodes:
                translated = translated_map.get(unit.id)
                if translated and translated.translated_text is not None:
                    elem.text = translated.translated_text
            part.tree.write(part.path, xml_declaration=True, encoding="utf-8", method="xml")

    def _repack_pptx(self, temp_root: Path, output_path: Path) -> None:
        if output_path.exists():
            output_path.unlink()
        with zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
            for folder, _, files in os.walk(temp_root):
                for filename in files:
                    abs_path = Path(folder) / filename
                    arcname = abs_path.relative_to(temp_root)
                    # Ensure POSIX separators inside the zip
                    zf.write(abs_path, arcname=str(arcname).replace(os.sep, "/"))

    def _discover_xml_parts(self, temp_root: Path) -> List[Path]:
        """
        Collect relevant XML parts that may contain text.
        """
        candidates: List[Path] = []
        slides_dir = temp_root / "ppt" / "slides"
        notes_dir = temp_root / "ppt" / "notesSlides"
        masters_dir = temp_root / "ppt" / "slideMasters"
        layouts_dir = temp_root / "ppt" / "slideLayouts"

        candidates.extend(sorted(slides_dir.glob("slide*.xml")))
        if self.include_notes:
            candidates.extend(sorted(notes_dir.glob("notesSlide*.xml")))
        if self.include_masters:
            candidates.extend(sorted(masters_dir.glob("slideMaster*.xml")))
            candidates.extend(sorted(layouts_dir.glob("slideLayout*.xml")))

        existing = [path for path in candidates if path.exists()]
        return existing

    def _build_deck_profile(self, pptx_path: Path) -> DeckProfile:
        """
        Build a basic deck profile (titles, sections, keywords).
        """
        try:
            prs = Presentation(pptx_path)
        except Exception:
            return DeckProfile(title=None, section_titles=[], keywords=[], summary=None)

        title = None
        section_titles: List[str] = []
        texts: List[str] = []
        for slide in prs.slides:
            if slide.shapes.title and slide.shapes.title.text:
                if title is None:
                    title = slide.shapes.title.text.strip()
                section_titles.append(slide.shapes.title.text.strip())
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)

        keywords = _extract_top_keywords(texts, top_n=20)
        # Lightweight summary: join first few titles/bullets
        summary = "; ".join(section_titles[:5]) if section_titles else None
        return DeckProfile(title=title, section_titles=section_titles, keywords=keywords, summary=summary)

    def _combine_context(self, user_context: Optional[str], deck_profile: DeckProfile) -> str:
        sections = []
        if user_context:
            sections.append(user_context)
        profile_text = deck_profile.to_context_string()
        if profile_text:
            sections.append(profile_text)
        return "\n\n".join(sections)

    def _generate_glossary(self, path: Path, units: List[TranslatableUnit], deck_profile: DeckProfile) -> None:
        """
        Create a simple glossary suggestion file (CSV) based on frequent terms.
        """
        terms = _extract_top_keywords([u.source_text for u in units], top_n=50)
        path.parent.mkdir(parents=True, exist_ok=True)
        with path.open("w", encoding="utf-8") as f:
            f.write("source,target,notes\n")
            for term in terms:
                f.write(f"{term},,\n")
        if deck_profile.keywords:
            self.logger.info("Glossary seeded with %d frequent terms", len(terms))

    def _extract_image_units(self, pptx_path: Path) -> Tuple[List[TranslatableUnit], List[OcrTextRegion]]:
        if not self.ocr_backend:
            # attempt to instantiate pytesseract by default
            try:
                self.ocr_backend = PytesseractOcrBackend()
            except Exception as exc:  # pragma: no cover - optional dependency
                self.logger.warning("OCR backend unavailable: %s", exc)
                return [], []

        prs = Presentation(pptx_path)
        image_inputs: List[OcrImageInput] = []
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                image = shape.image
                image_inputs.append(
                    OcrImageInput(
                        slide_index=slide_idx,
                        shape_index=shape_idx,
                        image_name=image.filename or f"image{slide_idx}_{shape_idx}",
                        image_bytes=image.blob,
                        width_px=image.size[0],
                        height_px=image.size[1],
                    )
                )
        if not image_inputs:
            return [], []

        ocr_regions = self.ocr_backend.recognize(image_inputs, config=self.ocr_config)
        ocr_units: List[TranslatableUnit] = []
        for region in ocr_regions:
            unit_id = self._next_id()
            region.unit_id = unit_id
            ocr_units.append(
                TranslatableUnit(
                    id=unit_id,
                    location=f"slide{region.slide_index}_img{region.shape_index}_bbox",
                    source_text=region.source_text,
                    context="image_text",
                )
            )
            region.translated_text = None  # will be filled after translation
        return ocr_units, ocr_regions

    def _apply_image_overlays(
        self,
        pptx_path: Path,
        regions: List[OcrTextRegion],
        translated_map: Dict[str, TranslatableUnit],
    ) -> None:
        prs = Presentation(pptx_path)
        emu = lambda v: Emu(int(v))
        for region in regions:
            translated = None
            if region.unit_id and region.unit_id in translated_map:
                translated = translated_map[region.unit_id].translated_text or translated_map[region.unit_id].source_text
            if not translated:
                translated = region.source_text
            try:
                slide = prs.slides[region.slide_index]
                shape = slide.shapes[region.shape_index]
            except IndexError:
                continue
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            img_width_px, img_height_px = shape.image.size
            if img_width_px == 0 or img_height_px == 0:
                continue
            left_px, top_px, w_px, h_px = region.bbox
            scale_x = shape.width / img_width_px
            scale_y = shape.height / img_height_px
            left = shape.left + emu(left_px * scale_x)
            top = shape.top + emu(top_px * scale_y)
            width = emu(max(w_px * scale_x, shape.width * 0.25))
            height = emu(max(h_px * scale_y, shape.height * 0.15))

            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text_frame.text = translated
            fill = textbox.fill
            fill.solid()
            try:
                if shape.fill.type and shape.fill.fore_color:
                    fill.fore_color.rgb = shape.fill.fore_color.rgb
            except Exception:
                pass
            textbox.text_frame.word_wrap = True
            textbox.name = f"pptx_translate_overlay_{region.image_name}"
        prs.save(pptx_path)

    def _embed_run_metadata(self, pptx_path: Path, metadata: Dict[str, str]) -> None:
        with TemporaryDirectory() as tmp_dir:
            tmp_root = Path(tmp_dir)
            with zipfile.ZipFile(pptx_path, "r") as zf:
                zf.extractall(tmp_root)
            props_path = tmp_root / "docProps" / "custom.xml"
            props_path.parent.mkdir(parents=True, exist_ok=True)
            if props_path.exists():
                tree = ET.parse(props_path)
                root = tree.getroot()
            else:
                root = ET.Element("{%s}Properties" % NAMESPACES["cp"])
                tree = ET.ElementTree(root)

            existing_names = {prop.get("name") for prop in root.findall("cp:property", NAMESPACES)}
            pid_start = 2
            for prop in root.findall("cp:property", NAMESPACES):
                pid = prop.get("pid")
                if pid and pid.isdigit():
                    pid_start = max(pid_start, int(pid) + 1)

            for name, value in metadata.items():
                if name in existing_names:
                    # update
                    prop = root.find(f"cp:property[@name='{name}']", NAMESPACES)
                    if prop is not None and prop.find("vt:lpwstr", NAMESPACES) is not None:
                        prop.find("vt:lpwstr", NAMESPACES).text = value
                    continue
                prop = ET.SubElement(
                    root,
                    "{%s}property" % NAMESPACES["cp"],
                    {
                        "fmtid": "{D5CDD505-2E9C-101B-9397-08002B2CF9AE}",
                        "pid": str(pid_start),
                        "name": name,
                    },
                )
                pid_start += 1
                lpwstr = ET.SubElement(prop, "{%s}lpwstr" % NAMESPACES["vt"])
                lpwstr.text = value

            tree.write(props_path, xml_declaration=True, encoding="UTF-8")

            # Repack
            with zipfile.ZipFile(pptx_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
                for folder, _, files in os.walk(tmp_root):
                    for filename in files:
                        abs_path = Path(folder) / filename
                        arcname = abs_path.relative_to(tmp_root)
                        zf.write(abs_path, arcname=str(arcname).replace(os.sep, "/"))

    def _generate_qa_report(self, pptx_path: Path, translated_units: List[TranslatableUnit], run_id: str) -> None:
        issues: List[dict] = []
        for unit in translated_units:
            if not unit.translated_text:
                continue
            ratio = len(unit.translated_text) / max(len(unit.source_text), 1)
            if ratio >= self.qa_threshold_length_ratio:
                issues.append(
                    {
                        "location": unit.location,
                        "source_preview": unit.source_text[:120],
                        "translated_preview": unit.translated_text[:120],
                        "reason": f"length_ratio>{self.qa_threshold_length_ratio}",
                        "ratio": ratio,
                    }
                )

        report = {
            "run_id": run_id,
            "issues_count": len(issues),
            "issues": issues,
        }
        if not self.qa_report_path:
            return
        self.qa_report_path.parent.mkdir(parents=True, exist_ok=True)
        if self.qa_report_format == "markdown":
            lines = [f"# QA Report (run_id={run_id})", "", f"Issues: {len(issues)}", ""]
            for idx, issue in enumerate(issues, start=1):
                lines.append(f"{idx}. {issue['location']} - {issue['reason']} (ratio={issue['ratio']:.2f})")
            self.qa_report_path.write_text("\n".join(lines), encoding="utf-8")
        else:
            self.qa_report_path.write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")

    def _next_id(self) -> str:
        self._id_counter += 1
        return f"t{self._id_counter}"


def sanitize_output_path(
    input_path: Path,
    user_output: Path | None,
    target_lang: str,
    run_id: Optional[str] = None,
    no_run_id: bool = False,
) -> Path:
    """
    Produce a default output path if the user didn't supply one.
    Includes run_id in the filename unless disabled.
    """
    base = user_output or input_path.with_name(f"{input_path.stem}.{target_lang}.pptx")
    if no_run_id or not run_id:
        return base
    if base.suffix.lower() == ".pptx":
        return base.with_name(f"{base.stem}.{run_id}.pptx")
    return Path(str(base) + f".{run_id}")


def generate_run_id() -> str:
    ts = datetime.utcnow().strftime("%Y%m%d-%H%M%S")
    suffix = secrets.token_hex(2)
    return f"{ts}-{suffix}"


def _extract_top_keywords(texts: Iterable[str], top_n: int = 20) -> List[str]:
    from collections import Counter

    counter: Counter[str] = Counter()
    for text in texts:
        for token in _tokenize(text):
            counter[token.lower()] += 1
    most_common = [token for token, _ in counter.most_common(top_n)]
    return most_common


def _tokenize(text: str) -> List[str]:
    clean = []
    buffer: List[str] = []
    for ch in text:
        if ch.isalnum():
            buffer.append(ch)
        else:
            if buffer:
                clean.append("".join(buffer))
                buffer = []
    if buffer:
        clean.append("".join(buffer))
    return [t for t in clean if len(t) > 2]
