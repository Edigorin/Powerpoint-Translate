from __future__ import annotations

import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_translate.backends import DummyBackend
from pptx_translate.translator import NAMESPACES, PptxTranslator


def _create_sample_pptx(tmp_path: Path) -> Path:
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Hello Title"
    slide.placeholders[1].text = "Body text"

    textbox = slide.shapes.add_textbox(left=Inches(1), top=Inches(4), width=Inches(3), height=Inches(1))
    textbox.text_frame.text = "Extra box"

    slide.notes_slide.notes_text_frame.text = "Note content"

    output = tmp_path / "input.pptx"
    prs.save(output)
    return output


def _collect_texts(pptx_path: Path, include_notes: bool) -> list[str]:
    texts: list[str] = []
    with zipfile.ZipFile(pptx_path, "r") as zf:
        def add_from_part(name: str) -> None:
            with zf.open(name) as f:
                tree = ET.parse(f)
                root = tree.getroot()
                for elem in root.findall(".//a:t", namespaces=NAMESPACES):
                    if elem.text:
                        texts.append(elem.text)

        def names(prefix: str) -> list[str]:
            return sorted([n for n in zf.namelist() if n.startswith(prefix) and n.endswith(".xml")])

        for part in names("ppt/slides/slide"):
            add_from_part(part)
        if include_notes:
            for part in names("ppt/notesSlides/notesSlide"):
                add_from_part(part)
        for part in names("ppt/slideMasters/slideMaster"):
            add_from_part(part)
        for part in names("ppt/slideLayouts/slideLayout"):
            add_from_part(part)
    return texts


def _shape_counts(pptx_path: Path) -> list[int]:
    prs = Presentation(pptx_path)
    return [len(slide.shapes) for slide in prs.slides]


@pytest.mark.parametrize("include_notes", [True, False])
def test_round_trip_preserves_layout_and_translates_text(tmp_path: Path, include_notes: bool) -> None:
    input_path = _create_sample_pptx(tmp_path)
    output_path = tmp_path / "output.pptx"

    translator = PptxTranslator(
        backend=DummyBackend(),
        include_notes=include_notes,
        include_masters=True,
        dry_run=False,
    )

    original_texts = _collect_texts(input_path, include_notes)
    translated_units = translator.translate_file(
        input_path=input_path,
        output_path=output_path,
        source_lang="en",
        target_lang="de",
    )
    translated_texts = _collect_texts(output_path, include_notes)

    assert _shape_counts(input_path) == _shape_counts(output_path)
    assert len(original_texts) == len(translated_texts)
    assert len(original_texts) == len(translated_units)

    for src, tgt in zip(original_texts, translated_texts):
        assert tgt == f"[de] {src}"

    if include_notes:
        out_prs = Presentation(output_path)
        note_text = out_prs.slides[0].notes_slide.notes_text_frame.text
        assert "[de]" in note_text
